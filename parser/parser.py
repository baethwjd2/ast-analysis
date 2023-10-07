import sys
import zipfile
import os
import datetime
import xml.etree.ElementTree as ET
import zipfile
import xml.dom.minidom
from os import rename
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

# def analyze(file_path):
#     if file_path.endswith('.docx'):
#         document = Document(file_path)
#         text = ""
#         for paragraph in document.paragraphs:
#             text += paragraph.text+"\n"
#         print("[파일 내용]")
#         print(text)

#     elif file_path.endswith('.pptx'):
#         presentation = Presentation(file_path)
#         text = ""
#         for slide in presentation.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape,"text"):
#                     text += shape.text + "\n"
#         print("[파일 내용]")
#         print(text)

#     elif file_path.endswith('.xlsx'):
#         workbook = load_workbook(filename=file_path)
#         sheet = workbook.active
#         data=[]
#         for row in sheet.iter_rows(values_only=True):
#             data.append(row)
#         print("[파일 내용]")
#         print(row)
#     else:
#         print("지원하지 않는 파일 형식입니다.")

def convert(file_path):
    # 확장자별 zip 변환
    if file_path.endswith('.docx'):
        zip_file_path = file_path.replace('.docx', '.zip')
        try:
            rename(file_path, zip_file_path)
            analyze_zip_central_directory(zip_file_path)
            return zip_file_path
        except FileNotFoundError:
            print(f"파일 '{file_path}'을 찾을 수 없습니다.")
        except Exception as e:
            print(f"파일 변환 중 오류가 발생했습니다: {e}")

    elif file_path.endswith('.pptx'):
        zip_file_path = file_path.replace('.pptx', '.zip')
        try:
            rename(file_path, zip_file_path)
            analyze_zip_central_directory(zip_file_path)
            return zip_file_path
        except FileNotFoundError:
            print(f"파일 '{file_path}'을 찾을 수 없습니다.")
        except Exception as e:
            print(f"파일 변환 중 오류가 발생했습니다: {e}")

    elif file_path.endswith('.xlsx'):
        zip_file_path = file_path.replace('.xlsx', '.zip')
        try:
            rename(file_path, zip_file_path)
            analyze_zip_central_directory(zip_file_path)
            return zip_file_path
        except FileNotFoundError:
            print(f"파일 '{file_path}'을 찾을 수 없습니다.")
        except Exception as e:
            print(f"파일 변환 중 오류가 발생했습니다: {e}")
    else:
        print("지원하지 않는 파일 형식입니다.")

def add_prefix(input_string, zip_path):
    #prefix = "C:\\workspace\\parser\\"  # 백슬래시(\)는 이스케이프 문자로 두 번 써야 합니다.
    result_string = zip_path + "\\" + input_string
    return result_string

def print_xml(xml_path):
    try:
        tree = ET.parse(xml_path)
        root = tree.getroot()
        xml_content = ET.tostring(root, encoding='utf-8').decode('utf-8')
        print(xml_content)
    except ET.ParseError as e:
        print(f"XML 파싱 오류 : {e}")

def xml_from_zip(zip_path, user_input):
    try:
        with zipfile.ZipFile(zip_path, 'r') as zip_file:
            with zip_file.open(user_input) as xml_file:
                xml_content = xml_file.read().decode('utf-8')
                return xml_content
    except FileNotFoundError:
        print(f"파일 '{zip_path}'을 찾을 수 없습니다.")
    except Exception as e:
        print(f"파일을 읽는 동안 오류가 발생했습니다: {e}")


def analyze_zip_central_directory(zip_file_path):
    # ZIP 파일 열기
    with zipfile.ZipFile(zip_file_path, 'r') as zip_file:
        # 모든 파일에 대한 Central Directory 정보 출력
        for info in zip_file.infolist():
            print("File Name:", info.filename)
            print("Version Made By:", info.create_version)
            print("Version Needed To Extract:", info.extract_version)
            print("General Purpose Bit Flags:", hex(info.flag_bits))
            print("Compression Method:", info.compress_type)
            file_timestamp = os.path.getmtime(zip_file_path)
            last_mod_time = datetime.datetime.fromtimestamp(file_timestamp)
            print("Last Mod Time:", last_mod_time)
            print("CRC32 Checksum:", hex(info.CRC))
            print("Compressed Size:", info.compress_size)
            print("Uncompressed Size:", info.file_size)
            print("File Comment:", info.comment)
            print("Internal File Attributes:", hex(info.internal_attr))
            print("External File Attributes:", hex(info.external_attr))
            print("Local Header Offset:", hex(info.header_offset))
            print()

if __name__ == "__main__":
    # analyze(file_path)
    file_path = input("분석할 파일 경로 :")
    zip_path = convert(file_path)

    while True:
        user_input = input("내용을 확인하고 싶은 파일명(xml만 가능, exit은 종료) :")
        if user_input =="exit":
            break
        else:
            xml_content = xml_from_zip(zip_path, user_input)
            dom=xml.dom.minidom.parseString(xml_content)
            pretty_xml = dom.toprettyxml()
            print(pretty_xml)
            print()
    #xml_path = add_prefix(user_input, zip_path)
    #print_xml(xml_path)
